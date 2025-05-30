import datetime
import os
from io import BytesIO

from PIL import Image
from PIL.Image import isImageType
from django.conf import settings
from django.utils.translation import gettext as _
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

SEEDSOURCE_TITLE = getattr(settings, "SEEDSOURCE_TITLE", _("Seedlot Selection Tool"))


class PPTCreator(object):
    def __init__(self):
        self.presentation = None
        self.width = None
        self.height = None

    def degree_sign(self, s):
        return s.replace("&deg;", "°")

    def add_text(self, text_frame, lines):
        for line in lines:
            paragraph = text_frame.add_paragraph()

            for segment in line:
                text, size, bold = segment

                run = paragraph.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold

    def get_transfer_method_text(self, method, center):
        if method != "seedzone":
            method_text = _(
                "Custom transfer limits, climatic center based on the selected location"
            )
        elif center == "zone":
            method_text = _("Transfer limits and climatic center based on seed zone")
        else:
            method_text = _(
                "Transfer limits based on seed zone, climatic center based on the selected location"
            )

        return method_text

    def replace_shape_image(self, shape, image, slide):
        im_bytes = BytesIO()
        image.save(im_bytes, "PNG")
        new_shape = slide.shapes.add_picture(
            im_bytes, Inches(0.5), Inches(0.5), Inches(9), Inches(6)
        )
        old_pic = shape._element
        new_pic = new_shape._element
        old_pic.addnext(new_pic)
        old_pic.getparent().remove(old_pic)

    def replace_shape_text(self, shape, text):
        paragraph = shape.text_frame.paragraphs[0]
        for run in paragraph.runs[1:]:
            paragraph._p.remove(run._r)
        paragraph.runs[0].text = text

    def add_title_text(self, slide, title):
        shape = slide.shapes.add_textbox(
            Inches(0.41), Inches(0.23), Inches(9.18), Inches(0.5)
        )
        tf = shape.text_frame
        tf.text = title
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(24)
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    def render_template(self, context):
        for slide in self.presentation.slides:
            self.render_template_slide(slide, context)

    def render_template_slide(self, slide, context):
        for shape in slide.shapes:
            if shape.name not in context:
                continue

            value = context[shape.name]

            if callable(value):
                value(shape)

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if not isImageType(value):
                    raise TypeError(
                        "Template value {} must be an Image type".format(shape.name)
                    )

                self.replace_shape_image(shape, value, slide)

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if not isinstance(value, str):
                    raise TypeError(
                        "Template value {} must be a string".format(shape.name)
                    )

                self.replace_shape_text(shape, value)

    def add_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])

        # Delete placeholders
        for placeholder in slide.placeholders:
            placeholder.element.getparent().remove(placeholder.element)

        return slide

    def create_overview_slide(self, context):
        objective = context["objective"]
        location_label = context["location_label"]
        point = context["point"]
        elevation = context["elevation"]
        seedlot_year = context["seedlot_year"]
        site_year = context["site_year"]
        site_model = context["site_model"]
        method = context["method"]
        center = context["center"]
        location = (point["y"], point["x"])
        data_url = (
            "http://cfcg.forestry.ubc.ca/projects/climate-data/climatebcwna/#ClimateWNA"
        )
        method_text = self.get_transfer_method_text(method, center)

        slide = self.add_slide()
        self.add_title_text(
            slide,
            "{} - {}".format(
                SEEDSOURCE_TITLE, datetime.datetime.today().strftime("%m/%d/%Y")
            ),
        )

        # Body
        shape = slide.shapes.add_textbox(
            Inches(0.65), Inches(0.73), Inches(8.69), Inches(6.19)
        )
        shape.text_frame.word_wrap = True
        self.add_text(
            shape.text_frame,
            (
                ((_("Objective:") + " ", 18, True), (objective, 18, False)),
                (("", 18, False),),
                (
                    ("{}: ".format(location_label), 18, True),
                    ("{}, {}".format(*location), 18, False),
                ),
                (
                    (_("Elevation:") + " ", 18, True),
                    (_("{elevation} ft").format(elevation=elevation), 18, False),
                ),
                (("", 18, False),),
                ((_("Climate scenarios"), 24, True),),
                ((_("Seedlot climate:") + " ", 18, True), (seedlot_year, 18, False)),
                (
                    (_("Planting site climate: ") + " ", 18, True),
                    (" ".join((site_year, site_model or "")), 18, False),
                ),
                (("", 18, False),),
                (
                    (_("Transfer limit method:") + " ", 18, True),
                    (method_text, 18, False),
                ),
                (("\n", 18, False),),
                ((_("Data URL:") + " ", 12, True), (data_url, 12, False)),
            ),
        )

        # Hyperlink URL
        shape.text_frame.paragraphs[-1].runs[-1].hyperlink.address = data_url

    def create_variables_slide(self, variables):
        slide = self.add_slide()
        self.add_title_text(slide, _("Climate Variables"))

        num_rows = len(variables) + 1
        table = slide.shapes.add_table(
            num_rows,
            3,
            Inches(0.47),
            Inches(0.73),
            Inches(9.05),
            Inches(0.4) * num_rows,
        ).table

        cols = table.columns
        cols[0].width = Inches(4.59)
        cols[1].width = Inches(2.06)
        cols[2].width = Inches(2.4)

        # Headers
        table.cell(0, 0).text = _("Variable")
        table.cell(0, 1).text = _("Center")
        table.cell(0, 2).text = _("Transfer limit") + " (+/-)"

        for i, variable in enumerate(variables, start=1):
            units = self.degree_sign(variable["units"])
            center_label = " ".join((variable["value"], units))
            limit_label = "{} {}{}".format(
                variable["limit"],
                units,
                " ({})".format(_("modified")) if variable["modified"] else "",
            )

            table.cell(i, 0).text = variable["label"]
            table.cell(i, 1).text = center_label
            table.cell(i, 2).text = limit_label

    def create_custom_functions_slide(self, custom_functions):
        slide = self.add_slide()
        self.add_title_text(slide, _("Custom Functions"))

        num_rows = len(custom_functions) + 1
        table = slide.shapes.add_table(
            num_rows,
            4,
            Inches(0.47),
            Inches(0.73),
            Inches(9.05),
            Inches(0.4) * num_rows,
        ).table

        cols = table.columns
        cols[0].width = Inches(2.1)
        cols[1].width = Inches(3.75)
        cols[2].width = Inches(0.6)
        cols[2].width = Inches(0.9)

        # Headers
        table.cell(0, 0).text = _("Name")
        table.cell(0, 1).text = _("Function")
        table.cell(0, 2).text = _("Center")
        table.cell(0, 3).text = _("Transfer limit") + " (+/-)"

        for i, custom_function in enumerate(custom_functions, start=1):
            table.cell(i, 0).text = custom_function["name"]
            table.cell(i, 1).text = custom_function["func"]
            table.cell(i, 2).text = str(custom_function["value"])
            table.cell(i, 3).text = str(custom_function["transfer"])

    def create_constraints_slide(self, constraints):
        slide = self.add_slide()
        self.add_title_text(slide, _("Constraints"))

        num_rows = len(constraints) + 1
        table = slide.shapes.add_table(
            num_rows,
            3,
            Inches(0.47),
            Inches(0.73),
            Inches(9.05),
            Inches(0.4) * num_rows,
        ).table

        cols = table.columns
        cols[0].width = Inches(4.59)
        cols[1].width = Inches(2.06)
        cols[2].width = Inches(2.4)

        # Headers
        table.cell(0, 0).text = _("Constraint")
        table.cell(0, 1).text = _("Value")
        table.cell(0, 2).text = "{} (+/-)".format(_("Range"))

        for i, constraint in enumerate(constraints, start=1):
            if constraint["type"] == "shapefile":
                table.cell(i, 0).text = constraint["label"]
                table.cell(i, 1)._tc.set("gridSpan", str(2))
                table.cell(i, 1).text = constraint["filename"]
            else:
                table.cell(i, 0).text = constraint["label"]
                table.cell(i, 1).text = constraint["value"]
                table.cell(i, 2).text = constraint["range"]

    def add_presenter_notes(self, slide, context):
        text_frame = slide.notes_slide.notes_text_frame

        objective = context["objective"]
        location_label = context["location_label"]
        point = context["point"]
        elevation = context["elevation"]
        seedlot_year = context["seedlot_year"]
        site_year = context["site_year"]
        site_model = context["site_model"]
        method = context["method"]
        center = context["center"]

        location = (point["y"], point["x"])
        method_text = self.get_transfer_method_text(method, center)

        lines = [
            ((_("Objective:") + " ", 12, True), (objective, 12, False)),
            (
                ("{}: ".format(location_label), 12, True),
                ("{}, {}".format(*location), 12, False),
            ),
            ((_("Elevation:") + " ", 12, True), ("{} ft".format(elevation), 12, False)),
            ((_("Climate Scenarios"), 12, True),),
            (
                ("  {} ".format(_("Seedlot climate:")), 12, True),
                (seedlot_year, 12, False),
            ),
            (
                ("  {} ".format(_("Planting site climate:")), 12, True),
                ("{} {}".format(site_year, site_model or ""), 12, False),
            ),
            ((_("Transfer limit method:") + " ", 12, True), (method_text, 12, False)),
        ]

        if method == "seedzone":
            band = context["band"]
            band_str = ", {}' - {}'".format(band[0], band[1]) if band else ""
            lines += [
                ((_("Species:") + " ", 12, True), (context["species"], 12, False)),
                (
                    (_("Seed zone:") + " ", 12, True),
                    (context["zone"] + band_str, 12, False),
                ),
            ]

        # Variables table
        variables = context["variables"]
        name_width = (
            max([len(_("Variable"))] + [len(x["label"]) for x in variables]) + 3
        )
        center_width = (
            max(
                [len(_("Center"))]
                + [
                    len(" ".join([str(x["value"]), self.degree_sign(x["units"])]))
                    for x in variables
                ]
            )
            + 3
        )
        transfer_width = max(
            [len(_("Transfer limit") + " (+/-)")]
            + [
                len(
                    "{} {}{}".format(
                        x["limit"],
                        self.degree_sign(x["units"]),
                        " ({})".format(_("modified")) if x["modified"] else "",
                    )
                )
                for x in variables
            ]
        )

        lines += [
            (("", 12, False),),
            ((_("Variables"), 12, True),),
            (
                (
                    "".join(
                        [
                            _("Variable").ljust(name_width),
                            _("Center").ljust(center_width),
                            _("Transfer limit") + " (+/-)".ljust(transfer_width),
                        ]
                    ),
                    12,
                    False,
                ),
            ),
            (("-" * (name_width + center_width + transfer_width), 12, False),),
        ]

        for variable in context["variables"]:
            units = self.degree_sign(variable["units"])
            lines += [
                (
                    (
                        "".join(
                            [
                                variable["label"].ljust(name_width),
                                "{} {}".format(variable["value"], units).ljust(
                                    center_width
                                ),
                                "{} {}{}".format(
                                    variable["limit"],
                                    units,
                                    (
                                        " ({})".format(_("modified"))
                                        if variable["modified"]
                                        else ""
                                    ),
                                ),
                            ]
                        ),
                        12,
                        False,
                    ),
                )
            ]

        if context["constraints"]:
            # Constraints table
            constraints = context["constraints"]
            name_width = (
                max([len("Constraint")] + [len(x["label"]) for x in constraints]) + 3
            )
            value_width = (
                max(
                    [len(_("Value"))]
                    + [
                        len(x["value"])
                        for x in [c for c in constraints if c["type"] != "shapefile"]
                    ]
                )
                + 3
            )
            range_width = (
                max(
                    [len(_("Range") + " (+/-)")]
                    + [
                        len(x["range"])
                        for x in [c for c in constraints if c["type"] != "shapefile"]
                    ]
                )
                + 3
            )

            # Ensure we have room for shapefile name, if there is one
            shape_constraint = [c for c in constraints if c["type"] == "shapefile"]
            if shape_constraint:
                filename_width = len(shape_constraint[0]["filename"])
                if filename_width > value_width + range_width:
                    range_width = filename_width - value_width

            lines += [
                (("", 12, False),),
                ((_("Constraints"), 12, True),),
                (
                    (
                        "".join(
                            [
                                _("Constraint").ljust(name_width),
                                _("Value").ljust(value_width),
                                _("Range") + " (+/-)".ljust(range_width),
                            ]
                        ),
                        12,
                        False,
                    ),
                ),
                (("-" * (name_width + value_width + range_width), 12, False),),
            ]

            for constraint in constraints:
                if constraint["type"] == "shapefile":
                    lines += [
                        (
                            (
                                "".join(
                                    [
                                        constraint["label"].ljust(name_width),
                                        constraint["filename"].ljust(
                                            value_width + range_width
                                        ),
                                    ]
                                ),
                                12,
                                False,
                            ),
                        )
                    ]
                else:
                    lines += [
                        (
                            (
                                "".join(
                                    [
                                        constraint["label"].ljust(name_width),
                                        constraint["value"].ljust(value_width),
                                        constraint["range"].ljust(range_width),
                                    ]
                                ),
                                12,
                                False,
                            ),
                        )
                    ]

        self.add_text(text_frame, lines)

        for paragraph in text_frame.paragraphs:
            paragraph.font.name = "Andale Mono"

    def get_presentation(self, context):
        self.presentation = Presentation(
            os.path.join(os.path.dirname(__file__), "templates", "pptx", "report.pptx")
        )
        self.width = Inches(self.presentation.slide_width / Inches(1))
        self.height = Inches(self.presentation.slide_height / Inches(1))

        self.render_template(
            dict(
                coord_bottom=self.degree_sign(context["south"]),
                coord_right=self.degree_sign(context["east"]),
                coord_left=self.degree_sign(context["west"]),
                coord_top=self.degree_sign(context["north"]),
                scale_label=context["scale"],
                map_image=Image.open(context["image_data"]),
                attribution=_("Generated {date} by the Seedlot Selection Tool").format(
                    date=datetime.datetime.today().strftime("%m/%d/%Y")
                ),
            )
        )
        self.create_overview_slide(context)
        self.create_variables_slide(context["variables"])

        if context["custom_functions"]:
            self.create_custom_functions_slide(context["custom_functions"])

        if context["constraints"]:
            self.create_constraints_slide(context["constraints"])

        self.add_presenter_notes(self.presentation.slides[0], context)

        return self.presentation
